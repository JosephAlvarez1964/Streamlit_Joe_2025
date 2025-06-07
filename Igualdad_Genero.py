from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Crear la presentación
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # diseño en blanco

# Función para agregar diapositivas
def add_slide(title, content):
    slide = prs.slides.add_slide(slide_layout)
    # Fondo azul claro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 229, 255)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = content
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0, 0, 0)

# Contenido de las diapositivas
slides_data = [
    ("Infografía Estatal: Chihuahua — Avances en Igualdad y Armonización Legislativa",
     "Panorama estatal en función de avances y armonizaciones legislativas según el marco normativo nacional y estatal."),
    ("Constitución Política para el Estado de Chihuahua",
     "• Armonizada con la Constitución Política de los Estados Unidos Mexicanos y tratados internacionales.\n"
     "• Reconoce la igualdad entre mujeres y hombres como principio fundamental."),
    ("Igualdad entre Mujeres y Hombres",
     "• Ley de Igualdad entre Mujeres y Hombres del Estado de Chihuahua (reforma 2024).\n"
     "• Elimina discriminación en ámbitos público y privado.\n"
     "• Establece acciones afirmativas e institucionales.\n"
     "• Principios: igualdad, equidad de género, no discriminación.\n"
     "• Reconoce derechos sin distinción.\n"
     "• Crea Unidades de Igualdad de Género.\n"
     "• Fomenta eliminación de estereotipos."),
    ("Vida Libre de Violencia para Mujeres y Niñas",
     "• Obligación del Estado: erradicar violencia y discriminación.\n"
     "• Implementar políticas públicas para prevenir y sancionar violencia.\n"
     "• El Consejo del Instituto Chihuahuense de las Mujeres define lineamientos."),
    ("Programas y Política Estatal para la Igualdad",
     "• Sistema Estatal para la Igualdad coordina acciones.\n"
     "• Programa Estatal para la Igualdad ejecuta y evalúa políticas.\n"
     "• Plan Estatal de Desarrollo incluye perspectiva de género.\n"
     "• Ley de Planeación aplica enfoque de género."),
    ("Avance y Armonización Legislativa",
     "• Chihuahua presenta 26 de 32 respuestas positivas.\n"
     "• Cuenta con instrumentos jurídicos y de planeación.\n"
     "• Se ubica en el rango más alto en transversalización de perspectiva de género."),
]

# Crear las diapositivas
for title, content in slides_data:
    add_slide(title, content)

# Guardar el archivo
prs.save("Infografia_Chihuahua_Igualdad.pptx")
